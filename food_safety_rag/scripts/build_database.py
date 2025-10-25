# # 1. File reading functions        # read_txt, read_pdf, read_docx, read_pptx
# import os
# from pypdf import PdfReader
# from docx import Document
# from pptx import Presentation

# def read_txt(file_path):
#     with open(file_path, "r", encoding="utf-8") as f:
#         return f.read()

# def read_pdf(file_path):
#     text = ""
#     reader = PdfReader(file_path)
#     for page in reader.pages:
#         text += page.extract_text() + "\n"
#     return text

# def read_docx(file_path):
#     doc = Document(file_path)
#     return "\n".join([para.text for para in doc.paragraphs])

# def read_pptx(file_path):
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text

# def read_file(file_path):
#     ext = file_path.split(".")[-1].lower()
#     if ext == "txt":
#         return read_txt(file_path)
#     elif ext == "pdf":
#         return read_pdf(file_path)
#     elif ext == "docx":
#         return read_docx(file_path)
#     elif ext == "pptx":
#         return read_pptx(file_path)
#     else:
#         return ""


# # 2. Hierarchical chunking
# from transformers import AutoTokenizer

# tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-mpnet-base-v2")

# def hierarchical_chunk(text):
#     """
#     Splits text into hierarchical chunks.
#     Returns list of dicts: { 'text', 'level', 'parent_id' }
#     """
#     chunks = []
#     level0_size, level1_size, level2_size = 2048, 512, 128
    
#     # Tokenize full text
#     tokens = tokenizer.encode(text, truncation=False)
    
#     chunk_id = 0
#     level0_chunks = [tokens[i:i+level0_size] for i in range(0, len(tokens), level0_size)]
    
#     for l0 in level0_chunks:
#         l0_text = tokenizer.decode(l0)
#         chunks.append({"text": l0_text, "level": 0, "parent_id": None, "chunk_id": chunk_id})
#         parent_id_l1 = chunk_id
#         chunk_id += 1
        
#         # Level 1
#         l1_chunks = [l0[i:i+level1_size] for i in range(0, len(l0), level1_size)]
#         for l1 in l1_chunks:
#             l1_text = tokenizer.decode(l1)
#             chunks.append({"text": l1_text, "level": 1, "parent_id": parent_id_l1, "chunk_id": chunk_id})
#             parent_id_l2 = chunk_id
#             chunk_id += 1
            
#             # Level 2
#             l2_chunks = [l1[i:i+level2_size] for i in range(0, len(l1), level2_size)]
#             for l2 in l2_chunks:
#                 l2_text = tokenizer.decode(l2)
#                 chunks.append({"text": l2_text, "level": 2, "parent_id": parent_id_l2, "chunk_id": chunk_id})
#                 chunk_id += 1
    
#     return chunks


# # 3. Elasticsearch index creation  # es = Elasticsearch(...), create index if not exists
# from elasticsearch import Elasticsearch
# import os
# from dotenv import load_dotenv

# load_dotenv()
# ES_URL = os.getenv("ES_URL")
# INDEX_NAME = os.getenv("INDEX_NAME")

# es = Elasticsearch(ES_URL)

# index_mapping = {
#     "mappings": {
#         "properties": {
#             "text": {"type": "text"},
#             "level": {"type": "integer"},
#             "parent_id": {"type": "integer"},
#             "source": {"type": "keyword"},
#             "embedding": {"type": "dense_vector", "dims": 768}  # depends on embedding model
#         }
#     }
# }

# if not es.indices.exists(index=INDEX_NAME):
#     es.indices.create(index=INDEX_NAME, body=index_mapping)
#     print(f"Created index {INDEX_NAME}")


# # 4. Compute Embeddings
# from sentence_transformers import SentenceTransformer
# import torch
# from tqdm import tqdm

# model = SentenceTransformer(os.getenv("MAIN_EMBED_MODEL"))

# def embed_texts(texts):
#     embeddings = model.encode(texts, convert_to_tensor=True, batch_size=16, show_progress_bar=True)
#     return embeddings


# # 5. Read Files and Prepare for Batch Processing
# import glob

# data_folder = os.path.join(os.path.dirname(__file__), '..', 'data') # More robust path
# all_files = glob.glob(os.path.join(data_folder, "*"))

# all_chunks = []
# all_texts_to_embed = []

# print("Step 1: Reading and chunking files...")
# for file_path in tqdm(all_files, desc="Processing files"):
#     if not os.path.isfile(file_path):
#         continue
        
#     try:
#         text = read_file(file_path)
#         if not text or not text.strip():
#             print(f"Warning: No text extracted from {file_path}. Skipping.")
#             continue
            
#         chunks = hierarchical_chunk(text)
        
#         # Add source info and prepare for embedding
#         for c in chunks:
#             c["source"] = os.path.basename(file_path)
#             all_chunks.append(c)
#             all_texts_to_embed.append(c["text"])
            
#     except Exception as e:
#         print(f"Error processing file {file_path}: {e}")

# print(f"Total chunks created: {len(all_chunks)}")


# # 6. Compute All Embeddings in a Single Batch
# print("\nStep 2: Computing embeddings for all chunks...")
# if all_texts_to_embed:
#     all_embeddings = embed_texts(all_texts_to_embed)
# else:
#     all_embeddings = []
#     print("No texts to embed. Exiting.")

# # 7. Batch Index to Elasticsearch
# print("\nStep 3: Indexing documents into Elasticsearch...")
# if all_chunks and len(all_embeddings) > 0:
#     from elasticsearch.helpers import bulk
    
#     actions = []
#     for chunk, embedding in tqdm(zip(all_chunks, all_embeddings), total=len(all_chunks), desc="Preparing for indexing"):
#         action = {
#             "_index": INDEX_NAME,
#             "_source": {
#                 "text": chunk["text"],
#                 "level": chunk["level"],
#                 "parent_id": chunk["parent_id"],
#                 "source": chunk["source"],
#                 "embedding": embedding.tolist()
#             }
#         }
#         actions.append(action)

#     # Perform the bulk indexing
#     try:
#         success, failed = bulk(es, actions, chunk_size=100, request_timeout=200)
#         print(f"Indexing complete! Success: {success}, Failed: {failed}")
#     except Exception as e:
#         print(f"Error during bulk indexing: {e}")
# else:
#     print("No data was indexed.")



# ==============================================================
# RAG KNOWLEDGE BASE BUILDER
# ==============================================================

import os
import glob
from tqdm import tqdm
from dotenv import load_dotenv
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from transformers import AutoTokenizer
from sentence_transformers import SentenceTransformer
from elasticsearch import Elasticsearch
from elasticsearch.helpers import bulk
import torch

# --------------------------------------------------------------
# 1. Utility: File Reading Functions
# --------------------------------------------------------------

def read_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def read_pdf(file_path):
    text = ""
    reader = PdfReader(file_path)
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"
    return text

def read_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_file(file_path):
    ext = file_path.split(".")[-1].lower()
    if ext == "txt":
        return read_txt(file_path)
    elif ext == "pdf":
        return read_pdf(file_path)
    elif ext == "docx":
        return read_docx(file_path)
    elif ext == "pptx":
        return read_pptx(file_path)
    else:
        return ""


# --------------------------------------------------------------
# 2. Hierarchical Chunking
# --------------------------------------------------------------

tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-mpnet-base-v2")

def hierarchical_chunk(text):
    """
    Splits text into hierarchical chunks.
    Returns list of dicts: { 'text', 'level', 'parent_id' }
    """
    chunks = []
    level0_size, level1_size, level2_size = 2048, 512, 128

    tokens = tokenizer.encode(text, truncation=False)
    chunk_id = 0
    level0_chunks = [tokens[i:i + level0_size] for i in range(0, len(tokens), level0_size)]

    for l0 in level0_chunks:
        l0_text = tokenizer.decode(l0)
        chunks.append({"text": l0_text, "level": 0, "parent_id": None, "chunk_id": chunk_id})
        parent_id_l1 = chunk_id
        chunk_id += 1

        # Level 1
        l1_chunks = [l0[i:i + level1_size] for i in range(0, len(l0), level1_size)]
        for l1 in l1_chunks:
            l1_text = tokenizer.decode(l1)
            chunks.append({"text": l1_text, "level": 1, "parent_id": parent_id_l1, "chunk_id": chunk_id})
            parent_id_l2 = chunk_id
            chunk_id += 1

            # Level 2
            l2_chunks = [l1[i:i + level2_size] for i in range(0, len(l1), level2_size)]
            for l2 in l2_chunks:
                l2_text = tokenizer.decode(l2)
                chunks.append({"text": l2_text, "level": 2, "parent_id": parent_id_l2, "chunk_id": chunk_id})
                chunk_id += 1

    return chunks


# --------------------------------------------------------------
# 3. Elasticsearch Index Creation
# --------------------------------------------------------------

load_dotenv()

ES_URL = os.getenv("ES_URL", "http://localhost:9200")
INDEX_NAME = os.getenv("INDEX_NAME", "food_safety_kb").lower()
MAIN_EMBED_MODEL = os.getenv("MAIN_EMBED_MODEL", "sentence-transformers/all-mpnet-base-v2")

# Initialize client
es = Elasticsearch(
    ES_URL,
    request_timeout=60,
    retry_on_timeout=True,
    max_retries=5
)

# Check connection
try:
    info = es.info()
    print(f"✅ Connected to Elasticsearch {info['version']['number']} at {ES_URL}")
except Exception as e:
    print("❌ Could not connect to Elasticsearch:", e)
    exit(1)

# Index mapping (ES 8.x compliant)
index_mapping = {
    "mappings": {
        "properties": {
            "text": {"type": "text"},
            "level": {"type": "integer"},
            "parent_id": {"type": "integer"},
            "source": {"type": "keyword"},
            "embedding": {
                "type": "dense_vector",
                "dims": 768,
                "index": False
            }
        }
    },
    "settings": {
        "index": {
            "number_of_shards": 1,
            "number_of_replicas": 0
        }
    }
}

try:
    if not es.indices.exists(index=INDEX_NAME):
        print(f"🆕 Creating index '{INDEX_NAME}' ...")
        es.indices.create(index=INDEX_NAME, body=index_mapping, ignore=400)
        print(f"✅ Index '{INDEX_NAME}' created successfully.")
    else:
        print(f"✅ Index '{INDEX_NAME}' already exists.")
except Exception as e:
    print(f"❌ Error creating index '{INDEX_NAME}': {e}")
    exit(1)


# --------------------------------------------------------------
# 4. Compute Embeddings
# --------------------------------------------------------------

print(f"\nLoading embedding model: {MAIN_EMBED_MODEL}")
model = SentenceTransformer(MAIN_EMBED_MODEL)

def embed_texts(texts):
    embeddings = model.encode(
        texts,
        convert_to_tensor=True,
        batch_size=16,
        show_progress_bar=True
    )
    return embeddings


# --------------------------------------------------------------
# 5. Read Files and Prepare for Batch Processing
# --------------------------------------------------------------

data_folder = os.path.join(os.path.dirname(__file__), '..', 'data')
all_files = glob.glob(os.path.join(data_folder, "*"))

all_chunks = []
all_texts_to_embed = []

print("\nStep 1: Reading and chunking files...")
for file_path in tqdm(all_files, desc="Processing files"):
    if not os.path.isfile(file_path):
        continue

    try:
        text = read_file(file_path)
        if not text or not text.strip():
            print(f"⚠️  Warning: No text extracted from {file_path}. Skipping.")
            continue

        chunks = hierarchical_chunk(text)

        for c in chunks:
            c["source"] = os.path.basename(file_path)
            all_chunks.append(c)
            all_texts_to_embed.append(c["text"])

    except Exception as e:
        print(f"❌ Error processing file {file_path}: {e}")

print(f"📄 Total chunks created: {len(all_chunks)}")


# --------------------------------------------------------------
# 6. Compute All Embeddings in a Single Batch
# --------------------------------------------------------------

print("\nStep 2: Computing embeddings for all chunks...")
if all_texts_to_embed:
    all_embeddings = embed_texts(all_texts_to_embed)
else:
    all_embeddings = []
    print("⚠️  No texts to embed. Exiting.")
    exit(0)


# --------------------------------------------------------------
# 7. Batch Index to Elasticsearch
# --------------------------------------------------------------

print("\nStep 3: Indexing documents into Elasticsearch...")

if all_chunks and len(all_embeddings) > 0:
    actions = []
    for chunk, embedding in tqdm(zip(all_chunks, all_embeddings), total=len(all_chunks), desc="Preparing for indexing"):
        action = {
            "_index": INDEX_NAME,
            "_source": {
                "text": chunk["text"],
                "level": chunk["level"],
                "parent_id": chunk["parent_id"],
                "source": chunk["source"],
                "embedding": embedding.tolist()
            }
        }
        actions.append(action)

    try:
        success, failed = bulk(es, actions, chunk_size=100, request_timeout=200)
        print(f"✅ Indexing complete! Success: {success}, Failed: {failed}")
    except Exception as e:
        print(f"❌ Error during bulk indexing: {e}")
else:
    print("⚠️  No data was indexed.")
