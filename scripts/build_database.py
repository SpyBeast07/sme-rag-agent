"""
FINAL BUILD_DATABASE.PY
Fully corrected, production-ready.
"""

import os
import glob
import pandas as pd
from tqdm import tqdm
from dotenv import load_dotenv

# File readers
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# Elasticsearch
from elasticsearch import Elasticsearch
from elasticsearch.helpers import bulk

# NLP
from transformers import AutoTokenizer
from sentence_transformers import SentenceTransformer
import re, hashlib

# ------------------------------------------------------------
# 1. Load environment and initialize Elasticsearch
# ------------------------------------------------------------
load_dotenv()

ES_URL = os.getenv("ES_URL")
INDEX_NAME = os.getenv("INDEX_NAME")
MAIN_EMBED_MODEL = os.getenv("MAIN_EMBED_MODEL")

print(f"📌 Using embedding model: {MAIN_EMBED_MODEL}")

es = Elasticsearch(ES_URL)

index_mapping = {
    "mappings": {
        "properties": {
            "text": {"type": "text"},
            "level": {"type": "integer"},
            "parent_id": {"type": "integer"},
            "source": {"type": "keyword"},
            "embedding": {"type": "dense_vector", "dims": 768},
            # Uncomment below if using domain model too
            # "embedding_domain": {"type": "dense_vector", "dims": 768},
            "metadata": {"type": "object"}
        }
    }
}

if not es.indices.exists(index=INDEX_NAME):
    es.indices.create(index=INDEX_NAME, body=index_mapping)
    print(f"✅ Created index: {INDEX_NAME}")
else:
    print(f"ℹ️ Index '{INDEX_NAME}' already exists.")


# ------------------------------------------------------------
# 2. Load metadata.csv
# ------------------------------------------------------------
metadata_df = pd.read_csv("reports/outputs/metadata.csv")
metadata_map = metadata_df.set_index("filename").to_dict(orient="index")

print(f"📌 Loaded metadata for {len(metadata_map)} files.")


# ------------------------------------------------------------
# 3. File reading utilities
# ------------------------------------------------------------
def read_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except:
        return ""

def read_pdf(path):
    text = ""
    try:
        reader = PdfReader(path)
        for page in reader.pages:
            t = page.extract_text() or ""
            text += t + "\n"
    except:
        pass
    return text

def read_docx(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except:
        return ""

def read_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except:
        pass
    return text

def read_file(path):
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".txt": return read_txt(path)
    if ext == ".pdf": return read_pdf(path)
    if ext == ".docx": return read_docx(path)
    if ext == ".pptx": return read_pptx(path)
    return ""


# ------------------------------------------------------------
# 4. Normalization + Dedupe helpers
# ------------------------------------------------------------
def normalize_text(t: str) -> str:
    t = t.replace("\r", "\n")
    t = re.sub(r"\n\s*\n+", "\n\n", t)
    t = re.sub(r"\s+", " ", t)
    return t.strip()

def chunk_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


# ------------------------------------------------------------
# 5. Paragraph-aware chunker (max 2048 tokens)
# ------------------------------------------------------------
tokenizer = AutoTokenizer.from_pretrained(MAIN_EMBED_MODEL)

def paragraph_chunker(text, max_tokens=2048, overlap=256):
    paragraphs = [p.strip() for p in re.split(r"\n{1,}", text) if p.strip()]
    
    chunks = []
    current = ""

    for p in paragraphs:
        candidate = (current + "\n\n" + p).strip() if current else p
        tok_len = len(tokenizer.encode(candidate))

        if tok_len <= max_tokens:
            current = candidate
        else:
            if current:
                chunks.append(current)

            # If paragraph is too large → split into sliding windows
            toks = tokenizer.encode(p)
            for i in range(0, len(toks), max_tokens - overlap):
                window = tokenizer.decode(toks[i:i+max_tokens])
                chunks.append(window)

            current = ""

    if current:
        chunks.append(current)

    return chunks


# ------------------------------------------------------------
# 6. Hierarchical chunking (2048 → 512 → 128 tokens)
# ------------------------------------------------------------
def hierarchical_chunk(text):
    final_chunks = []
    global_chunk_id = 0

    # --- Level 0: paragraph-based (2048 tokens) ---
    lvl0_chunks = paragraph_chunker(text, max_tokens=2048)

    for lvl0 in lvl0_chunks:
        lvl0_tokens = tokenizer.encode(lvl0)
        lvl0_id = global_chunk_id
        final_chunks.append({
            "text": lvl0,
            "level": 0,
            "parent_id": None,
            "chunk_id": lvl0_id
        })
        global_chunk_id += 1

        # --- Level 1: 512 tokens ---
        lvl1_token_chunks = [
            lvl0_tokens[i:i+512] for i in range(0, len(lvl0_tokens), 512)
        ]

        for lvl1_tok in lvl1_token_chunks:
            lvl1_text = tokenizer.decode(lvl1_tok)
            lvl1_id = global_chunk_id
            final_chunks.append({
                "text": lvl1_text,
                "level": 1,
                "parent_id": lvl0_id,
                "chunk_id": lvl1_id
            })
            global_chunk_id += 1

            # --- Level 2: 128 tokens ---
            lvl2_token_chunks = [
                lvl1_tok[i:i+128] for i in range(0, len(lvl1_tok), 128)
            ]
            for lvl2_tok in lvl2_token_chunks:
                lvl2_text = tokenizer.decode(lvl2_tok)
                final_chunks.append({
                    "text": lvl2_text,
                    "level": 2,
                    "parent_id": lvl1_id,
                    "chunk_id": global_chunk_id
                })
                global_chunk_id += 1

    return final_chunks


# ------------------------------------------------------------
# 7. Embedding Model
# ------------------------------------------------------------
embed_general = SentenceTransformer(MAIN_EMBED_MODEL)


# ------------------------------------------------------------
# 8. Process all documents
# ------------------------------------------------------------
data_folder = "data"
all_files = glob.glob(os.path.join(data_folder, "**/*.*"), recursive=True)

all_docs_chunks = []
all_texts = []

print("\n📘 Processing files...")
seen_hashes = set()

for path in tqdm(all_files):
    if not os.path.isfile(path):
        continue

    filename = os.path.basename(path)
    text = read_file(path)

    if not text.strip():
        continue

    raw_chunks = hierarchical_chunk(text)

    for c in raw_chunks:
        t = normalize_text(c["text"])
        if len(t) < 50:
            continue

        h = chunk_hash(t)
        if h in seen_hashes:
            continue
        seen_hashes.add(h)

        c["text"] = t
        c["filename"] = filename
        c["metadata"] = metadata_map.get(filename, {})

        all_docs_chunks.append(c)
        all_texts.append(t)

print(f"✅ Total final chunks: {len(all_docs_chunks)}")


# ------------------------------------------------------------
# 9. Embed and index
# ------------------------------------------------------------
import numpy as np
from elasticsearch.helpers import streaming_bulk

MODEL_MAX_TOKENS = 512  # mpnet / most SentenceTransformers models
BATCH_SIZE = 8          # lower if you get OOM; increase if OK
INDEX_BATCH_SIZE = 200  # number of docs per ES bulk

# Build lists of items to embed and index, but FILTER long chunks (>MODEL_MAX_TOKENS)
to_index = []
to_embed_texts = []
to_embed_indices = []  # track mapping from embedding result to to_index

def clean_metadata(meta):
    clean = {}
    for k, v in meta.items():
        if pd.isna(v):
            clean[k] = ""
        elif isinstance(v, (int, float)):
            clean[k] = str(v)
        else:
            clean[k] = str(v)
    return clean

print("\n🔎 Filtering chunks to ensure token length <= model max ...")
for idx, chunk in enumerate(all_docs_chunks):
    chunk["metadata"] = clean_metadata(chunk["metadata"])
    tok_len = len(tokenizer.encode(chunk["text"]))
    if tok_len > MODEL_MAX_TOKENS:
        # skip or optionally truncate here; we skip (you already have smaller level2 chunks)
        continue
    # append doc for embedding+index
    to_index.append(chunk)
    to_embed_texts.append(chunk["text"])
    to_embed_indices.append(len(to_index)-1)

print(f"➡️ Chunks to embed & index (<= {MODEL_MAX_TOKENS} tokens): {len(to_index)}")

if len(to_index) == 0:
    raise SystemExit("No chunks are small enough for embedding. Check your chunker.")

# Batch-encode with small batch size and convert to float32 numpy
print("\n🧠 Computing embeddings in batches...")
embeddings = []
for i in tqdm(range(0, len(to_embed_texts), BATCH_SIZE), desc="Embedding batches"):
    batch_texts = to_embed_texts[i:i+BATCH_SIZE]
    emb_batch = embed_general.encode(batch_texts, convert_to_numpy=True, normalize_embeddings=True, batch_size=BATCH_SIZE)
    # ensure float32 and extend
    if isinstance(emb_batch, np.ndarray):
        emb_batch = emb_batch.astype(np.float32)
        embeddings.extend(list(emb_batch))
    else:
        # fallback: convert tensors
        emb_batch = np.asarray(emb_batch).astype(np.float32)
        embeddings.extend(list(emb_batch))

print(f"✅ Computed {len(embeddings)} embeddings.")

# Build indexing actions in small chunks and use streaming_bulk to capture errors
print("\n📤 Indexing into Elasticsearch in batches (with error capture)...")
def gen_actions():
    for i, chunk in enumerate(to_index):
        emb = embeddings[i]
        doc = {
            "text": chunk["text"],
            "level": chunk["level"],
            "parent_id": chunk["parent_id"],
            "source": chunk["filename"],
            "embedding": emb.tolist(),
            "metadata": chunk["metadata"]
        }
        yield {
            "_index": INDEX_NAME,
            "_source": doc
        }

# Use streaming_bulk which yields (ok, result) pairs so we can log failures
failures = []
count_indexed = 0
for ok, item in streaming_bulk(client=es, actions=gen_actions(), chunk_size=INDEX_BATCH_SIZE, max_retries=3, request_timeout=200):
    if not ok:
        failures.append(item)
    else:
        count_indexed += 1

print(f"✅ Indexed documents: {count_indexed}")
if failures:
    print(f"⚠️ Number of failures: {len(failures)} — showing up to 5 examples:")
    for e in failures[:5]:
        print(e)
    # stop and advise to inspect failures
    raise SystemExit("Bulk indexing had failures. See printed failures above.")
